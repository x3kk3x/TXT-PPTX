import tkinter as tk
from tkinter import filedialog, ttk, messagebox
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

class txtPPTXApp:
    def __init__(self, root):
        self.root = root
        self.root.title("txtPPTX")
        self.root.geometry("400x400")

        # Variables
        self.folder_path = tk.StringVar(value="")
        self.output_path = tk.StringVar(value="")
        self.font = tk.StringVar(value="Calibri")
        self.theme = tk.StringVar(value="Minimalist")
        self.words_per_slide = tk.IntVar(value=100)
        self.include_images = tk.BooleanVar(value=False)
        self.status = tk.StringVar(value="Spremno za početak")
        self.default_folder = tk.BooleanVar(value=False)

        # UI Elements (Matching Mockup)
        tk.Label(root, text="Izaberi Fasciklu sa Tekstovima:").pack(pady=5)
        tk.Entry(root, textvariable=self.folder_path, width=50).pack(pady=5)
        tk.Button(root, text="Pregledaj", command=self.browse_folder).pack(pady=5)
        tk.Checkbutton(root, text="Sačuvaj kao podrazumevanu fasciklu", variable=self.default_folder).pack(pady=5)

        self.settings_frame = tk.Frame(root)
        tk.Button(self.settings_frame, text="Prikaži Napredna Podešavanja", command=self.toggle_settings).pack()
        tk.Label(self.settings_frame, text="Font:").pack()
        tk.OptionMenu(self.settings_frame, self.font, "Calibri", "Arial", "Roboto").pack()
        tk.Label(self.settings_frame, text="Tema:").pack()
        tk.OptionMenu(self.settings_frame, self.theme, "Minimalist", "Office", "Plavi Gradijent").pack()
        tk.Label(self.settings_frame, text="Reči po Slajdu:").pack()
        tk.Scale(self.settings_frame, from_=50, to=150, orient="horizontal", variable=self.words_per_slide).pack()
        tk.Checkbutton(self.settings_frame, text="Uključi Slike", variable=self.include_images).pack()
        self.settings_frame.pack(pady=5, fill="x")
        self.settings_frame.pack_forget()

        tk.Button(root, text="Kreiraj Prezentacije", command=self.generate_presentations).pack(pady=10)

        tk.Label(root, textvariable=self.status).pack(pady=5)
        self.progress = ttk.Progressbar(root, length=300, mode="determinate")
        self.progress.pack(pady=5)

        tk.Label(root, text="Sačuvaj Prezentacije U:").pack(pady=5)
        tk.Entry(root, textvariable=self.output_path, width=50).pack(pady=5)
        tk.Button(root, text="Pregledaj", command=self.browse_output_folder).pack(pady=5)

        tk.Label(root, text="v1.0").pack(side="left", padx=10, pady=5)
        tk.Button(root, text="Pomoć", command=self.show_help).pack(side="right", padx=10, pady=5)

        self.settings_visible = False

    def browse_folder(self):
        folder = filedialog.askdirectory()
        if folder:
            self.folder_path.set(folder)
            if not self.output_path.get():
                self.output_path.set(os.path.join(folder, "output"))

    def browse_output_folder(self):
        folder = filedialog.askdirectory()
        if folder:
            self.output_path.set(folder)

    def toggle_settings(self):
        if self.settings_visible:
            self.settings_frame.pack_forget()
            self.settings_visible = False
        else:
            self.settings_frame.pack(pady=5, fill="x")
            self.settings_visible = True

    def show_help(self):
        messagebox.showinfo("Pomoć", "1. Izaberi fasciklu sa .txt fajlovima.\n2. Podesi font, temu i reči po slajdu (opciono).\n3. Klikni 'Kreiraj Prezentacije'.\n4. Pronađi .pptx fajlove u izlaznoj fascikli.")

    def generate_presentations(self):
        folder = self.folder_path.get()
        output_dir = self.output_path.get()
        if not folder or not os.path.exists(folder):
            messagebox.showerror("Greška", "Izaberi validnu fasciklu sa tekstovima!")
            return
        if not output_dir:
            messagebox.showerror("Greška", "Izaberi izlaznu fasciklu!")
            return

        os.makedirs(output_dir, exist_ok=True)
        txt_files = [f for f in os.listdir(folder) if f.endswith(".txt")]
        if not txt_files:
            messagebox.showerror("Greška", "Nema .txt fajlova u fascikli!")
            return

        self.progress["maximum"] = len(txt_files)
        self.status.set("Obrada u toku...")
        self.root.update()

        for i, file_name in enumerate(txt_files):
            self.create_pptx(os.path.join(folder, file_name), output_dir)
            self.progress["value"] = i + 1
            self.status.set(f"Obrada {i + 1} od {len(txt_files)} fajlova")
            self.root.update()

        self.status.set(f"Uspešno kreirano {len(txt_files)} prezentacija!")
        messagebox.showinfo("Uspeh", "Prezentacije su kreirane u izlaznoj fascikli!")

    def create_pptx(self, file_path, output_dir):
        prs = Presentation()

        # Title Slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = os.path.splitext(os.path.basename(file_path))[0].replace("_", " ")
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title.text_frame.paragraphs[0].font.name = self.font.get()
        title.text_frame.paragraphs[0].font.size = Pt(36)

        # Content Slides
        slide_layout = prs.slide_layouts[1]
        current_slide = None
        word_count = 0

        with open(file_path, 'r', encoding='utf-8') as f:
            lines = f.read().splitlines()
            for line in lines:
                if word_count > self.words_per_slide.get():
                    current_slide = None
                    word_count = 0

                if not current_slide:
                    current_slide = prs.slides.add_slide(slide_layout)
                    content_box = current_slide.shapes.placeholders[1].text_frame
                    content_box.clear()

                if line.strip():
                    if line.startswith("# "):
                        title = current_slide.shapes.title
                        title.text = line[2:].strip()
                        title.text_frame.paragraphs[0].font.name = self.font.get()
                        title.text_frame.paragraphs[0].font.size = Pt(28)
                    else:
                        p = content_box.add_paragraph()
                        p.text = line.strip()
                        p.font.name = self.font.get()
                        p.font.size = Pt(20)
                        if line.startswith("- "):
                            p.level = 1
                            p.text = line[2:].strip()
                        word_count += len(line.split())

        # Save Presentation
        output_file = os.path.join(output_dir, f"{os.path.splitext(os.path.basename(file_path))[0]}_Prezentacija.pptx")
        prs.save(output_file)

if __name__ == "__main__":
    root = tk.Tk()
    app = txtPPTXApp(root)
    root.mainloop()